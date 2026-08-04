from pathlib import Path

from pptx import Presentation


BASE = Path(__file__).resolve().parent
SRC = BASE / "discussion_concept_deck_left_text_style_v3.pptx"
OUT = BASE / "discussion_concept_deck_left_text_style_v4_no_text_5_8_11.pptx"
TARGET_SLIDES = {5, 8, 11}


def remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def main():
    prs = Presentation(str(SRC))
    for index, slide in enumerate(prs.slides, start=1):
        if index not in TARGET_SLIDES:
            continue
        for shape in list(slide.shapes):
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                remove_shape(shape)
    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    main()
