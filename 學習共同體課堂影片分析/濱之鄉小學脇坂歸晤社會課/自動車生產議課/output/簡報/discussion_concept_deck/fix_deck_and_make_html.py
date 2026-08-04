import html
import json
import re
from pathlib import Path

from pptx import Presentation


BASE = Path(__file__).resolve().parent
PPTX = BASE / "discussion_concept_deck_editable.pptx"
FIXED_PPTX = BASE / "discussion_concept_deck_editable_fixed.pptx"
SPEC = BASE / "deck_spec.json"
HTML_OUT = BASE / "discussion_concept_deck_responsive.html"


def remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def is_question_garbage(text):
    compact = re.sub(r"\s+", "", text or "")
    return bool(compact) and len(compact) >= 5 and set(compact) <= {"?"}


def fix_paragraph_text(text):
    text = text or ""
    # Replace mojibake bullet markers at the start of body lines with a stable symbol.
    return re.sub(r"^([ \t]*)([?？])([ \t]+)", r"\1▸\3", text)


def fix_pptx():
    prs = Presentation(str(PPTX))

    for slide_index, slide in enumerate(prs.slides, start=1):
        shapes_to_remove = []

        for shape in list(slide.shapes):
            if slide_index == 8 and getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                # The approved image already contains the slide 8 text; remove the duplicate editable right-side text.
                if shape.left > prs.slide_width * 0.45:
                    shapes_to_remove.append(shape)
                    continue

            if not (getattr(shape, "has_text_frame", False) and shape.has_text_frame):
                continue

            full_text = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
            if is_question_garbage(full_text):
                shapes_to_remove.append(shape)
                continue

            for paragraph in shape.text_frame.paragraphs:
                if is_question_garbage(paragraph.text):
                    paragraph.clear()
                    continue
                paragraph.text = fix_paragraph_text(paragraph.text)

        for shape in shapes_to_remove:
            remove_shape(shape)

        if slide_index == 8:
            # Remove the right-side panel and accent rule that framed the duplicate text.
            for shape in list(slide.shapes):
                if not getattr(shape, "has_text_frame", False) and shape.left > prs.slide_width * 0.45:
                    # Keep the background picture; remove generated panel/rule shapes only.
                    if getattr(shape, "shape_type", None) != 13:
                        remove_shape(shape)

    prs.save(str(FIXED_PPTX))
    return FIXED_PPTX


def css():
    return """
:root {
  color-scheme: light;
  --ink: #342f27;
  --muted: #6f6658;
  --paper: #f6eedc;
  --panel: rgba(248, 240, 222, 0.93);
  --accent: #af5337;
}
* { box-sizing: border-box; }
html { scroll-behavior: smooth; }
body {
  margin: 0;
  font-family: "Microsoft JhengHei", "Noto Sans TC", system-ui, sans-serif;
  background: #efe3ca;
  color: var(--ink);
}
.deck-nav {
  position: sticky;
  top: 0;
  z-index: 10;
  display: flex;
  gap: .45rem;
  overflow-x: auto;
  padding: .6rem .8rem;
  background: rgba(246, 238, 220, .92);
  border-bottom: 1px solid rgba(52,47,39,.12);
  backdrop-filter: blur(8px);
}
.deck-nav a {
  flex: 0 0 auto;
  color: var(--ink);
  text-decoration: none;
  font-size: .86rem;
  padding: .28rem .5rem;
  border: 1px solid rgba(52,47,39,.16);
  border-radius: 999px;
}
.slide {
  min-height: 100svh;
  display: grid;
  place-items: center;
  padding: clamp(12px, 2vw, 28px);
  border-bottom: 1px solid rgba(52,47,39,.1);
}
.canvas {
  position: relative;
  width: min(100%, 1280px);
  aspect-ratio: 16 / 9;
  overflow: hidden;
  background: var(--paper);
  box-shadow: 0 16px 38px rgba(45, 38, 28, .16);
}
.canvas img {
  position: absolute;
  inset: 0;
  width: 100%;
  height: 100%;
  object-fit: cover;
}
.panel {
  position: absolute;
  left: var(--x);
  top: var(--y);
  width: var(--w);
  min-height: var(--h);
  padding: clamp(18px, 2.5vw, 34px);
  background: var(--panel);
  border-radius: min(4vw, 42px);
}
.panel h1,
.panel h2 {
  margin: 0;
  line-height: 1.12;
  letter-spacing: 0;
  color: var(--ink);
}
.panel h1 { font-size: clamp(2.1rem, 4.4vw, 4.1rem); }
.panel h2 { font-size: clamp(1.45rem, 2.75vw, 2.45rem); }
.rule {
  width: min(7rem, 28%);
  height: 4px;
  margin: 1rem 0 1.1rem;
  background: var(--accent);
}
ul {
  list-style: none;
  margin: 0;
  padding: 0;
  display: grid;
  gap: .48rem;
}
li {
  position: relative;
  padding-left: 1.35em;
  font-size: clamp(1rem, 1.72vw, 1.45rem);
  line-height: 1.42;
}
li::before {
  content: "▸";
  position: absolute;
  left: 0;
  color: var(--accent);
  font-weight: 700;
}
.source {
  margin-top: 1rem;
  color: var(--muted);
  font-size: clamp(.78rem, 1.05vw, 1rem);
  line-height: 1.35;
}
.slide-8 .panel { display: none; }
@media (max-width: 760px) {
  .slide {
    min-height: auto;
    padding: 10px;
  }
  .canvas {
    aspect-ratio: auto;
    min-height: 100svh;
  }
  .canvas img {
    position: relative;
    height: auto;
    display: block;
  }
  .panel {
    position: relative;
    left: auto;
    top: auto;
    width: auto;
    min-height: 0;
    margin: -10vw 10px 12px;
  }
  .slide-8 .panel { display: none; }
}
"""


PANEL_MAP = {
    1: ("4%", "9%", "47%", "68%"),
    2: ("4%", "7%", "41%", "76%"),
    3: ("4%", "7%", "40%", "76%"),
    4: ("4%", "7%", "41%", "76%"),
    5: ("3%", "6%", "40%", "80%"),
    6: ("4%", "7%", "41%", "76%"),
    7: ("4%", "7%", "42%", "76%"),
    8: ("54%", "6%", "42%", "80%"),
    9: ("3.5%", "7%", "42%", "76%"),
    10: ("3.5%", "7%", "43%", "76%"),
    11: ("3.5%", "6%", "43%", "80%"),
    12: ("3.5%", "7%", "43%", "76%"),
    13: ("3.5%", "7%", "44%", "76%"),
    14: ("4%", "9%", "44%", "68%"),
}


SOURCE_NOTES = {
    1: "來源：中文議課逐字稿；日文逐字稿與中日比較僅作概念校正",
    2: "來源：議課開場與觀課者回到學生學習事實的發言",
    3: "來源：議課中對社會科本質、真實兩難與關係探究的討論",
    4: "來源：約 00:16:00 起，家長訪談與汽車產業生活化設計",
    5: "來源：議課中對「一欄表」、學生筆記與不發言學生的分析",
    6: "來源：議課中對國內生產、海外生產與企業存續兩難的討論",
    7: "來源：議課中江次、藤田、內田等小組思考轉折的分析",
    8: "來源：觀課教師與校長對沉默、等待、廣賴與桑島的分析",
    9: "來源：議課中對佐佐木指名、眼神與參與保障的討論",
    10: "來源：議課中對新聞事件、海外生產與全球市民素養的討論",
    11: "來源：觀課者發言中專業觀看方式的轉變",
    12: "來源：佐藤學評講對課例與學習共同體哲學的連結",
    13: "來源：中文逐字稿概念整理；日文逐字稿與比較報告校正",
    14: "來源：議課整體收束與學習共同體哲學概念整理",
}


def make_html():
    data = json.loads(SPEC.read_text(encoding="utf-8"))
    nav = []
    sections = []
    for slide in data["slides"]:
        number = int(slide["number"])
        title = slide["title"]
        nav.append(f'<a href="#s{number:02d}">{number:02d}</a>')
        x, y, w, h = PANEL_MAP[number]
        title_tag = "h1" if number == 1 else "h2"
        bullets = "\n".join(f"<li>{html.escape(point)}</li>" for point in slide.get("key_points", []))
        panel = "" if number == 8 else f"""
      <article class="panel" style="--x:{x};--y:{y};--w:{w};--h:{h};">
        <{title_tag}>{html.escape(title)}</{title_tag}>
        <div class="rule"></div>
        <ul>
          {bullets}
        </ul>
        <p class="source">{html.escape(SOURCE_NOTES.get(number, ""))}</p>
      </article>"""
        sections.append(f"""
  <section id="s{number:02d}" class="slide slide-{number}">
    <div class="canvas">
      <img src="origin_image/slide_{number:02d}.png" alt="{html.escape(title)}">
      {panel}
    </div>
  </section>""")

    doc = f"""<!doctype html>
<html lang="zh-Hant">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>{html.escape(data.get("goal", "學習共同體議課簡報"))}</title>
  <style>{css()}</style>
</head>
<body>
  <nav class="deck-nav" aria-label="投影片導覽">
    {''.join(nav)}
  </nav>
  {''.join(sections)}
</body>
</html>
"""
    HTML_OUT.write_text(doc, encoding="utf-8")
    return HTML_OUT


if __name__ == "__main__":
    print(fix_pptx())
    print(make_html())
