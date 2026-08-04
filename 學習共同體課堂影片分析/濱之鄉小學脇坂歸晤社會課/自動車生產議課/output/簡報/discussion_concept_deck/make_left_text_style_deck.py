import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


BASE = Path(__file__).resolve().parent
SPEC = BASE / "deck_spec.json"
OUT = BASE / "discussion_concept_deck_left_text_style_v3.pptx"


INK = RGBColor(38, 35, 30)
PAPER = RGBColor(248, 240, 222)


def set_font(run, size, bold=False):
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = INK


def add_textbox(slide, x, y, w, h, text, size, bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.space_after = Pt(0)
    p.line_spacing = 1.08
    r = p.add_run()
    r.text = text
    set_font(r, size, bold)
    return box


def add_soft_paper_wash(slide, x=0, y=0, w=5.9, h=4.6, transparency=6):
    # A flat paper wash, not the old rounded text card. It cleans baked-in raster text on a few slides.
    wash = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    wash.fill.solid()
    wash.fill.fore_color.rgb = PAPER
    wash.fill.transparency = transparency
    wash.line.fill.background()


def add_slide(prs, slide_spec):
    number = int(slide_spec["number"])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    image_path = BASE / "origin_image" / f"slide_{number:02d}.png"
    if number == 8 and (BASE / "slide_08_clean_for_left_text.png").exists():
        image_path = BASE / "slide_08_clean_for_left_text.png"
    slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Keep the screenshot style: no panel, direct upper-left type.
    # Slides with baked-in generated text need a faint paper wash so editable text is legible.
    if number in {5, 11}:
        add_soft_paper_wash(slide)
    if number == 8:
        add_soft_paper_wash(slide, 0, 0, 5.9, 4.6, 6)

    title = slide_spec["title"]
    points = slide_spec.get("key_points", [])

    title_h = 0.8 if len(title) <= 17 else 1.2
    add_textbox(slide, 0.62, 0.78, 5.25, title_h, title, 28, True)

    y = 1.68 if title_h < 1.0 else 2.08
    for point in points[:4]:
        lines = max(1, (len(point) + 16) // 17)
        h = 0.34 * lines + 0.12
        add_textbox(slide, 0.32, y, 5.45, h, "▸ " + point, 20, False)
        y += h + 0.24


def main():
    spec = json.loads(SPEC.read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    for slide_spec in spec["slides"]:
        add_slide(prs, slide_spec)

    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    main()
