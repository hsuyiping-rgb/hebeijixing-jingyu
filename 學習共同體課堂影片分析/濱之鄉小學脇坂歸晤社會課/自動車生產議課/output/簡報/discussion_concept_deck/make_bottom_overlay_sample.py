import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


BASE = Path(__file__).resolve().parent
SPEC = BASE / "deck_spec.json"
OUT = BASE / "discussion_concept_deck_bottom_overlay_sample.pptx"


def set_font(run, size, bold=False, color=RGBColor(52, 47, 39)):
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, size, bold=False, color=RGBColor(52, 47, 39)):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.space_after = Pt(0)
    r = p.add_run()
    r.text = text
    set_font(r, size, bold, color)
    return box


def main():
    spec = json.loads(SPEC.read_text(encoding="utf-8"))
    slide_spec = spec["slides"][0]

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.shapes.add_picture(str(BASE / "origin_image" / "slide_01.png"), 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Bottom reading band: visually acts as a base edge, but is not the old large rounded text box.
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.18), prs.slide_width, Inches(2.32))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(248, 240, 222)
    band.fill.transparency = 13
    band.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.18), prs.slide_width, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(175, 83, 55)
    accent.line.fill.background()

    add_text(slide, 0.62, 5.38, 4.2, 0.45, slide_spec["title"], 28, True)

    points = slide_spec["key_points"]
    left_points = points[:2]
    right_points = points[2:]
    for i, point in enumerate(left_points):
        add_text(slide, 0.66, 5.98 + i * 0.45, 5.55, 0.38, "▸ " + point, 18)
    for i, point in enumerate(right_points):
        add_text(slide, 6.78, 5.98 + i * 0.45, 5.8, 0.38, "▸ " + point, 18)

    add_text(slide, 0.66, 7.02, 10.5, 0.24, "來源：中文議課逐字稿；日文逐字稿與中日比較僅作概念校正", 11, False, RGBColor(103, 94, 80))

    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    main()
