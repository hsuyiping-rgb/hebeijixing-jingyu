from pathlib import Path
import shutil
from PIL import Image

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "簡報" / "output" / "notebooklm_lesson_study_illustrated_pdf_3editable_revised_illustrations.pptx"
OUT = ROOT / "簡報" / "output" / "notebooklm_lesson_study_illustrated_pdf_3editable_safe_people.pptx"
ASSET_DIR = ROOT / "圖片" / "output" / "safe_people_illustrations"
ASSET_DIR.mkdir(parents=True, exist_ok=True)

GEN = Path.home() / ".codex" / "generated_images" / "019f3719-fb2a-74a2-b87d-62da698368a0"
RAW = {
    "full_class": GEN / "call_5dEeHPnusbJu99mVG66UJo47.png",
    "student_board": GEN / "call_ocdKgjfbpkemyU6E9SOom3fR.png",
    "group_graph": GEN / "call_8eOSrqsz0fFvqzlZrncDKfhW.png",
    "teacher_board": GEN / "call_tuql0spuTIkG6J7Qvkp6mdGW.png",
    "teacher_group": GEN / "call_a9I4wz9hpKxUgjyZOhrBHnlD.png",
}


def make_padded(src: Path, out: Path, margin_ratio: float = 0.03):
    if not src.exists():
        raise FileNotFoundError(src)
    im = Image.open(src).convert("RGB")
    canvas_w, canvas_h = 1920, 1080
    margin_x = int(canvas_w * margin_ratio)
    margin_y = int(canvas_h * margin_ratio)
    max_w = canvas_w - 2 * margin_x
    max_h = canvas_h - 2 * margin_y
    scale = min(max_w / im.width, max_h / im.height)
    new_w = int(im.width * scale)
    new_h = int(im.height * scale)
    resized = im.resize((new_w, new_h), Image.Resampling.LANCZOS)
    canvas = Image.new("RGB", (canvas_w, canvas_h), (244, 247, 241))
    canvas.paste(resized, ((canvas_w - new_w) // 2, (canvas_h - new_h) // 2))
    canvas.save(out)
    return out


def delete_shape(shape):
    el = shape._element
    el.getparent().remove(el)


def image_size_emu(path: Path):
    im = Image.open(path)
    return im.width, im.height


def add_contained_picture(slide, image_path: Path, left_in, top_in, width_in, height_in):
    left = Inches(left_in)
    top = Inches(top_in)
    box_w = Inches(width_in)
    box_h = Inches(height_in)
    img_w, img_h = image_size_emu(image_path)
    scale = min(box_w / img_w, box_h / img_h)
    pic_w = int(img_w * scale)
    pic_h = int(img_h * scale)
    pic_left = int(left + (box_w - pic_w) / 2)
    pic_top = int(top + (box_h - pic_h) / 2)
    return slide.shapes.add_picture(str(image_path), pic_left, pic_top, width=pic_w, height=pic_h)


def replace_picture(slide, image_path: Path, left=6.15, top=1.18, width=6.72, height=5.72):
    for sh in list(slide.shapes):
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            delete_shape(sh)
    add_contained_picture(slide, image_path, left, top, width, height)


def main():
    safe = {}
    for key, path in RAW.items():
        copied = ASSET_DIR / f"{key}_raw.png"
        padded = ASSET_DIR / f"{key}_safe_margin.png"
        shutil.copy2(path, copied)
        make_padded(copied, padded)
        safe[key] = padded

    prs = Presentation(str(SRC))

    # Slide numbers are 1-based.
    mapping = {
        1: ("full_class", (6.35, 1.22, 6.25, 5.50)),
        5: ("group_graph", (6.15, 1.18, 6.72, 5.72)),
        6: ("student_board", (6.15, 1.18, 6.72, 5.72)),
        7: ("group_graph", (6.15, 1.18, 6.72, 5.72)),
        8: ("student_board", (6.15, 1.18, 6.72, 5.72)),
        9: ("teacher_board", (6.15, 1.18, 6.72, 5.72)),
        10: ("teacher_group", (6.15, 1.18, 6.72, 5.72)),
        11: ("student_board", (6.15, 1.18, 6.72, 5.72)),
        12: ("full_class", (6.15, 1.18, 6.72, 5.72)),
        13: ("full_class", (6.15, 1.18, 6.72, 5.72)),
        14: ("group_graph", (6.15, 1.18, 6.72, 5.72)),
        15: ("group_graph", (6.15, 1.18, 6.72, 5.72)),
        16: ("group_graph", (6.15, 1.18, 6.72, 5.72)),
        17: ("full_class", (6.15, 1.18, 6.72, 5.72)),
        18: ("student_board", (6.15, 1.18, 6.72, 5.72)),
    }

    for slide_no, (key, box) in mapping.items():
        replace_picture(prs.slides[slide_no - 1], safe[key], *box)

    prs.save(str(OUT))
    print(OUT)
    print(len(prs.slides))


if __name__ == "__main__":
    main()
