from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "output" / "notebooklm_lesson_study.pptx"
DEST = ROOT / "output" / "notebooklm_lesson_study_illustrated.pptx"
IMAGE_DIR = ROOT / "output" / "images"
IMAGE_MAP = [1, 8, 12, 9, 6, 3, 4, 7, 2, 5, 9, 8, 11, 13, 12]


def add_cropped_picture(slide, path: Path, left, top, width, height):
    with Image.open(path) as image:
        ratio = image.width / image.height
    box_ratio = width / height
    picture = slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    if ratio > box_ratio:
        crop = (1 - box_ratio / ratio) / 2
        picture.crop_left = picture.crop_right = crop
    elif ratio < box_ratio:
        crop = (1 - ratio / box_ratio) / 2
        picture.crop_top = picture.crop_bottom = crop


deck = Presentation(SOURCE)
if len(deck.slides) != len(IMAGE_MAP):
    raise RuntimeError(f"Slide count mismatch: {len(deck.slides)}")

for index, (slide, image_number) in enumerate(zip(deck.slides, IMAGE_MAP), 1):
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape._element.getparent().remove(shape._element)

    text_shapes = [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
    title = text_shapes[0] if text_shapes else None
    body = text_shapes[1] if len(text_shapes) > 1 else None

    if title:
        title.left, title.top = Inches(0.65), Inches(0.34)
        title.width, title.height = Inches(12.0), Inches(0.78)
    if body:
        body.left, body.top = Inches(0.7), Inches(1.38)
        body.width, body.height = Inches(5.45), Inches(5.45)
        body.text_frame.word_wrap = True
        body.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        for paragraph in body.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(17.5 if index > 1 else 18.5)

    add_cropped_picture(
        slide,
        IMAGE_DIR / f"slide_{image_number:02d}.png",
        Inches(6.4), Inches(1.38), Inches(6.25), Inches(5.28),
    )

deck.save(DEST)
print(DEST)
