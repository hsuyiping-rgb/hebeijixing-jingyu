from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "output" / "slides.pptx"
DEST = ROOT / "output" / "slides_illustrated.pptx"
IMAGE_DIR = ROOT / "output" / "images"


def add_cropped_picture(slide, image_path: Path, left, top, width, height):
    with Image.open(image_path) as image:
        image_ratio = image.width / image.height
    box_ratio = width / height
    picture = slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    if image_ratio > box_ratio:
        visible = box_ratio / image_ratio
        crop = (1 - visible) / 2
        picture.crop_left = crop
        picture.crop_right = crop
    elif image_ratio < box_ratio:
        visible = image_ratio / box_ratio
        crop = (1 - visible) / 2
        picture.crop_top = crop
        picture.crop_bottom = crop
    return picture


presentation = Presentation(SOURCE)

for index, slide in enumerate(presentation.slides, 1):
    image_path = IMAGE_DIR / f"slide_{index:02d}.png"
    if not image_path.exists():
        raise FileNotFoundError(image_path)

    text_shapes = [shape for shape in slide.shapes if hasattr(shape, "text_frame") and shape.has_text_frame]
    title = text_shapes[0] if text_shapes else None
    body = text_shapes[1] if len(text_shapes) > 1 else None

    if title:
        title.left = Inches(0.65)
        title.top = Inches(0.35)
        title.width = Inches(12.05)
        title.height = Inches(0.75)

    if body:
        body.left = Inches(0.7)
        body.top = Inches(1.35)
        body.width = Inches(5.45)
        body.height = Inches(5.55)
        body.text_frame.word_wrap = True
        body.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        for paragraph in body.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(20 if index == 1 else 18)

    add_cropped_picture(
        slide,
        image_path,
        Inches(6.4),
        Inches(1.35),
        Inches(6.3),
        Inches(5.35),
    )

try:
    presentation.save(DEST)
    print(DEST)
except PermissionError:
    fallback = DEST.with_name("slides_illustrated_updated.pptx")
    presentation.save(fallback)
    print(fallback)
