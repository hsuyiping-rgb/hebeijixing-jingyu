from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("output")
SOURCE = OUT / "notebooklm_lesson_study_illustrated.pptx"
TARGET = OUT / "notebooklm_lesson_study_illustrated_pdf_3editable_slides.pptx"

prs = Presentation(str(SOURCE))
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

GREEN = RGBColor(35, 65, 45)
BG = RGBColor(239, 244, 238)
TEXT = RGBColor(76, 84, 78)
MUTED = RGBColor(105, 113, 106)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(202, 211, 199)
GRID = RGBColor(222, 224, 216)
BLUE = RGBColor(78, 112, 152)
RED = RGBColor(176, 86, 72)
ORANGE = RGBColor(202, 137, 52)
DARK = RGBColor(48, 58, 52)
NOTE = RGBColor(242, 248, 241)
NOTE_LINE = RGBColor(203, 219, 202)
PAPER = RGBColor(250, 248, 239)
PAPER_LINE = RGBColor(224, 217, 196)
FONT = "Microsoft JhengHei"


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color, width=1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_fill(bg, BG)
    bg.line.fill.background()


def add_text(slide, text, x, y, w, h, size=24, color=TEXT, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    for run in p.runs:
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_multiline(slide, text, x, y, w, h, size=23, color=TEXT, bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return box


def add_card(slide, x, y, w, h):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(shp, WHITE)
    set_line(shp, LINE, 1.5)
    return shp


def add_table(slide, x, y):
    rows, cols = 2, 5
    table_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(4.7), Inches(0.82))
    tbl = table_shape.table
    headers = ["時間(分)", "0", "5", "10", "15"]
    vals = ["離家距離(m)", "0", "500", "1000", "1500"]
    for r, row in enumerate([headers, vals]):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(250, 250, 247) if r == 0 or c == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(11 if c else 10)
                    run.font.color.rgb = TEXT
                    run.font.bold = r == 0 or c == 0
    for c in range(cols):
        tbl.columns[c].width = Inches(0.78 if c else 1.18)
    return table_shape


def draw_graph(slide, x, y, w, h, points, color, labels=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(rect, PAPER)
    set_line(rect, PAPER_LINE, 1)
    gx = x + 0.65
    gy = y + 0.45
    gw = w - 1.0
    gh = h - 0.9
    for i in range(29):
        xx = gx + gw * i / 28
        ln = slide.shapes.add_connector(1, Inches(xx), Inches(gy), Inches(xx), Inches(gy + gh))
        set_line(ln, GRID, 0.4)
    for j in range(16):
        yy = gy + gh - gh * j / 15
        ln = slide.shapes.add_connector(1, Inches(gx), Inches(yy), Inches(gx + gw), Inches(yy))
        set_line(ln, GRID, 0.4)
    for val in [0, 500, 1000, 1500]:
        yy = gy + gh - gh * val / 1500
        ln = slide.shapes.add_connector(1, Inches(gx), Inches(yy), Inches(gx + gw), Inches(yy))
        set_line(ln, RGBColor(115, 119, 111), 1)
        add_text(slide, str(val), x + 0.06, yy - 0.11, 0.7, 0.22, 10, TEXT)
    ln = slide.shapes.add_connector(1, Inches(gx), Inches(gy + gh), Inches(gx + gw), Inches(gy + gh))
    set_line(ln, DARK, 1.5)
    ln = slide.shapes.add_connector(1, Inches(gx), Inches(gy), Inches(gx), Inches(gy + gh))
    set_line(ln, DARK, 1.5)
    add_text(slide, "m", gx - 0.34, gy - 0.27, 0.3, 0.2, 10, TEXT)
    add_text(slide, "分", gx + gw - 0.25, gy + gh + 0.15, 0.3, 0.2, 10, TEXT)

    def pt(t, m):
        return gx + gw * t / 28, gy + gh - gh * m / 1500

    xy = [pt(t, m) for t, m in points]
    for (x1, y1), (x2, y2) in zip(xy, xy[1:]):
        ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        set_line(ln, color, 4)
    for xx, yy in xy:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(xx - 0.045), Inches(yy - 0.045), Inches(0.09), Inches(0.09))
        set_fill(dot, color)
        dot.line.fill.background()
    for t in [0, 5, 10, 15, 20, 25, 28]:
        xx = gx + gw * t / 28
        add_text(slide, str(t), xx - 0.08, gy + gh + 0.08, 0.35, 0.18, 9, TEXT)
    if labels:
        for text, t, m, c in labels:
            xx, yy = pt(t, m)
            add_text(slide, text, xx + 0.08, yy - 0.25, 1.2, 0.25, 12, c, bold=True)


def slide_pdf1():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_text(s, "PDF第1頁：與朋友的約定", 0.7, 0.45, 8.8, 0.55, 30, GREEN, True)
    add_text(s, "用表或折線圖表示「時間」與「離家距離」的關係。", 0.75, 1.1, 10.8, 0.35, 15, MUTED)
    add_card(s, 0.7, 1.7, 5.4, 5.4)
    add_text(s, "題目原文整理", 1.0, 2.02, 2.5, 0.38, 20, GREEN, True)
    add_multiline(
        s,
        "駿太跟朋友小智約好在離家 1500m 的公園見面。\n"
        "駿太預計用 1 分鐘走 100m 的速度移動。\n"
        "駿太出家門後，會在幾分鐘後到達公園呢？\n"
        "請用表或折線圖表示，進行思考。",
        1.0,
        2.55,
        4.75,
        1.8,
        17,
        TEXT,
    )
    rel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.55), Inches(4.7), Inches(1.2))
    set_fill(rel, NOTE)
    set_line(rel, NOTE_LINE, 1)
    add_text(s, "數學關係", 1.2, 4.75, 1.4, 0.3, 16, GREEN, True)
    add_multiline(s, "• 每 1 分鐘走 100m\n• 公園在離家 1500m 的位置\n• 問的是幾分鐘後到達", 1.2, 5.08, 4.1, 0.8, 13, TEXT)
    add_text(s, "表格整理", 1.0, 5.93, 1.6, 0.25, 15, GREEN, True)
    add_table(s, 1.0, 6.26)
    add_card(s, 6.45, 1.7, 6.1, 5.4)
    add_text(s, "折線圖表示", 6.85, 2.02, 2.6, 0.38, 20, GREEN, True)
    draw_graph(s, 6.85, 2.55, 5.25, 3.75, [(0, 0), (15, 1500)], BLUE, [("15分 / 1500m", 15, 1500, BLUE), ("100m/分", 6, 600, BLUE)])
    note = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(6.45), Inches(5.0), Inches(0.32))
    set_fill(note, NOTE)
    set_line(note, NOTE_LINE, 1)
    add_text(s, "固定速度形成一直線；到 1500m 時就是到達公園。", 7.15, 6.5, 4.8, 0.22, 12, GREEN)
    return s


def slide_pdf2():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_text(s, "PDF第2頁：實際行動的折線圖", 0.7, 0.45, 8.8, 0.55, 30, GREEN, True)
    add_text(s, "從折線圖判讀駿太實際怎麼行動。", 0.75, 1.1, 10.8, 0.35, 15, MUTED)
    add_card(s, 0.7, 1.65, 12.0, 5.6)
    add_multiline(s, "下面的折線圖表示了駿太的實際行動。\n從圖中，駿太是怎麼行動的呢？", 1.0, 2.0, 4.2, 1.0, 18, TEXT)
    draw_graph(s, 1.0, 3.05, 5.2, 3.35, [(0, 0), (5, 500), (10, 500), (15, 1500)], RED, [("停留", 6, 500, RED), ("加快前進", 12, 900, RED)])
    draw_graph(s, 6.85, 2.1, 5.3, 4.3, [(0, 0), (5, 500), (10, 500), (15, 1500)], RED, [("停留", 6, 500, RED), ("加快前進", 12, 900, RED)])
    note = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(6.55), Inches(11.1), Inches(0.35))
    set_fill(note, NOTE)
    set_line(note, NOTE_LINE, 1)
    add_text(s, "水平線表示時間增加，但離家的距離不變。", 1.2, 6.61, 10.5, 0.22, 13, GREEN)
    return s


def slide_pdf3():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_text(s, "PDF第3頁：等了2分鐘後的變化", 0.7, 0.45, 8.8, 0.55, 30, GREEN, True)
    add_text(s, "從後續折線圖判讀駿太的位置變化。", 0.75, 1.1, 10.8, 0.35, 15, MUTED)
    add_card(s, 0.7, 1.65, 12.0, 5.6)
    add_multiline(s, "駿太在時間內到達公園，並且等了 2 分鐘，小智都還沒到。\n下面的折線圖表示了駿太的什麼樣子呢？", 1.0, 1.95, 11.0, 0.9, 18, TEXT)
    draw_graph(
        s,
        1.0,
        3.05,
        11.1,
        3.4,
        [(0, 0), (5, 500), (10, 500), (15, 1500), (17, 1500), (20, 900), (21, 1200), (27, 0)],
        ORANGE,
        [("等2分鐘", 15.6, 1500, ORANGE), ("回到900m", 20, 900, ORANGE), ("又往公園方向", 21, 1200, ORANGE)],
    )
    note = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(6.58), Inches(11.1), Inches(0.35))
    set_fill(note, NOTE)
    set_line(note, NOTE_LINE, 1)
    add_text(s, "往下是靠近家；往上是靠近公園。未達 1500m 就沒有回到公園。", 1.2, 6.64, 10.5, 0.22, 13, GREEN)
    return s


new_slides = [slide_pdf1(), slide_pdf2(), slide_pdf3()]
sldIdLst = prs.slides._sldIdLst
ids = [sldIdLst[-3], sldIdLst[-2], sldIdLst[-1]]
for sid in ids:
    sldIdLst.remove(sid)
for offset, sid in enumerate(ids):
    sldIdLst.insert(1 + offset, sid)

prs.save(str(TARGET))
print(TARGET)
print(len(prs.slides))
