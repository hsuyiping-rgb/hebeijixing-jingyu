from pathlib import Path
import shutil

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "output" / "notebooklm_lesson_study_illustrated_pdf_3editable_slides.pptx"
OUT = ROOT / "output" / "notebooklm_lesson_study_illustrated_pdf_3editable_revised_illustrations.pptx"
IMG_DIR = ROOT / "output" / "revised_illustrations"
IMG_DIR.mkdir(parents=True, exist_ok=True)

GEN = Path.home() / ".codex" / "generated_images" / "019f3719-fb2a-74a2-b87d-62da698368a0"
SOURCE_IMAGES = {
    "panorama": GEN / "call_skuKf9l21L80KLJ9baxR6opd.png",
    "slide12_scene": ROOT / "tmp" / "extracted_slide_pictures" / "slide_12_pic_2.png",
    "group_graph": GEN / "call_YzrnQ8gWv2G7lIIDYtFSGVnp.png",
    "older_teacher_board": GEN / "call_eMPstjR7suQd0z9H8CzCRaKC.png",
    "older_teacher_group": GEN / "call_Lmacpg2ReFSvoBEb15EOM0q0.png",
    "pair_graph": GEN / "call_EIOqauLQo9N7IFtEJODfhAEl.png",
    "reflection": GEN / "call_e37fGYW3oigJ0IkI2KLzGUvt.png",
}

IMAGES = {}
for key, src in SOURCE_IMAGES.items():
    if not src.exists():
        raise FileNotFoundError(src)
    dst = IMG_DIR / f"{key}.png"
    shutil.copy2(src, dst)
    IMAGES[key] = dst


def delete_shape(shape):
    el = shape._element
    el.getparent().remove(el)


def replace_picture(slide, image_path, left=6.25, top=1.18, width=6.6, height=5.72):
    for sh in list(slide.shapes):
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            delete_shape(sh)
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), Inches(width), Inches(height))


def fit_picture(slide, image_path, left, top, width, height):
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), Inches(width), Inches(height))


def set_font_size(shape, size_pt, color=None):
    if not getattr(shape, "has_text_frame", False):
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size_pt)
            run.font.name = "Microsoft JhengHei"
            if color:
                run.font.color.rgb = color


def grow_all_text(prs):
    for slide in prs.slides:
        for sh in slide.shapes:
            if not getattr(sh, "has_text_frame", False):
                continue
            txt = sh.text.strip()
            if not txt:
                continue
            # Title-like boxes are near the top and wide.
            if sh.top < Inches(0.9) and sh.width > Inches(5):
                set_font_size(sh, 32, RGBColor(32, 75, 50))
            elif sh.width > Inches(4.5) and sh.height > Inches(2.5):
                set_font_size(sh, 18, RGBColor(77, 84, 78))
            else:
                # Small labels in graphs should remain readable but not oversized.
                set_font_size(sh, 11 if sh.height < Inches(0.35) else 15, RGBColor(77, 84, 78))


def remove_pdf_prefix_titles(prs):
    replacements = {
        "PDF第1頁：與朋友的約定": "與朋友的約定",
        "PDF第2頁：實際行動的折線圖": "實際行動的折線圖",
        "PDF第3頁：等了2分鐘後的變化": "等了2分鐘後的變化",
    }
    for slide in [prs.slides[1], prs.slides[2], prs.slides[3]]:
        for sh in slide.shapes:
            if getattr(sh, "has_text_frame", False):
                text = sh.text.strip()
                if text in replacements:
                    sh.text = replacements[text]
                    set_font_size(sh, 30, RGBColor(32, 75, 50))


def enlarge_slide3(prs):
    slide = prs.slides[2]

    # Delete the duplicated right graph and the old bottom note.
    for sh in list(slide.shapes):
        if sh.left > Inches(6.35) or (sh.top > Inches(6.35) and sh.left > Inches(0.8)):
            delete_shape(sh)

    # Enlarge question text.
    for sh in slide.shapes:
        if getattr(sh, "has_text_frame", False) and "下面的折線圖表示了駿太的實際行動" in sh.text:
            sh.left = Inches(1.0)
            sh.top = Inches(1.85)
            sh.width = Inches(10.8)
            sh.height = Inches(0.75)
            sh.text_frame.margin_left = Inches(0.05)
            set_font_size(sh, 21, RGBColor(77, 84, 78))

    # Move and scale the remaining graph cluster from left side into the center.
    # Original cluster bounds approx: x 1.0..6.2, y 3.05..6.4.
    old_left, old_top = Inches(1.0), Inches(3.05)
    new_left, new_top = Inches(1.35), Inches(2.75)
    sx, sy = 1.95, 1.13
    for sh in slide.shapes:
        if sh.left >= Inches(0.9) and sh.left < Inches(6.3) and sh.top >= Inches(3.0):
            sh.left = int(new_left + (sh.left - old_left) * sx)
            sh.top = int(new_top + (sh.top - old_top) * sy)
            sh.width = int(sh.width * sx)
            sh.height = int(sh.height * sy)
            if getattr(sh, "has_text_frame", False):
                set_font_size(sh, 12 if sh.height < Inches(0.35) else 15, RGBColor(77, 84, 78))

    # Add a new note with larger text.
    box = slide.shapes.add_shape(1, Inches(1.15), Inches(6.62), Inches(11.1), Inches(0.48))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(238, 244, 235)
    box.line.color.rgb = RGBColor(185, 204, 181)
    tx = slide.shapes.add_textbox(Inches(1.35), Inches(6.68), Inches(10.7), Inches(0.35))
    tx.text = "水平線表示時間增加，但離家的距離不變；斜率變大，表示移動速度變快。"
    tx.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    set_font_size(tx, 15, RGBColor(77, 84, 78))


def fix_slide2_layout(prs):
    slide = prs.slides[1]
    for sh in list(slide.shapes):
        if getattr(sh, "has_text_frame", False) and sh.text.strip() == "題目原文整理":
            delete_shape(sh)

    for sh in slide.shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        text = sh.text.strip()
        if text.startswith("駿太跟朋友小智約好"):
            sh.left = Inches(1.0)
            sh.top = Inches(2.08)
            sh.width = Inches(4.65)
            sh.height = Inches(1.55)
            set_font_size(sh, 15, RGBColor(77, 84, 78))
        elif text == "數學關係":
            sh.left = Inches(1.18)
            sh.top = Inches(4.12)
            sh.width = Inches(1.8)
            sh.height = Inches(0.25)
            set_font_size(sh, 12, RGBColor(77, 84, 78))
        elif "每 1 分鐘走 100m" in text:
            sh.text = "‧ 每 1 分鐘走 100m\n‧ 公園在離家 1500m 的位置"
            sh.left = Inches(1.18)
            sh.top = Inches(4.46)
            sh.width = Inches(4.3)
            sh.height = Inches(0.58)
            set_font_size(sh, 13, RGBColor(77, 84, 78))
        elif text == "表格整理":
            sh.left = Inches(1.0)
            sh.top = Inches(5.58)
            sh.width = Inches(1.8)
            sh.height = Inches(0.25)
            set_font_size(sh, 11, RGBColor(77, 84, 78))

    # Resize the green math relationship box so the bullets do not spill into the table.
    for sh in slide.shapes:
        if (
            sh.shape_type != MSO_SHAPE_TYPE.PICTURE
            and abs(sh.left - Inches(1.0)) < 10000
            and abs(sh.top - Inches(4.55)) < 10000
            and sh.width > Inches(4.0)
            and sh.height > Inches(0.8)
        ):
            sh.top = Inches(3.86)
            sh.height = Inches(1.25)

    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.TABLE:
            sh.left = Inches(1.0)
            sh.top = Inches(6.0)
            sh.width = Inches(4.3)
            sh.height = Inches(0.74)
            for row in sh.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        para.alignment = PP_ALIGN.CENTER
                        for run in para.runs:
                            run.font.size = Pt(9)
                            run.font.name = "Microsoft JhengHei"
                            run.font.color.rgb = RGBColor(77, 84, 78)


def main():
    prs = Presentation(str(SRC))
    grow_all_text(prs)
    remove_pdf_prefix_titles(prs)
    enlarge_slide3(prs)
    fix_slide2_layout(prs)

    # Replace illustrations. Slide numbering is 1-based in user request.
    replace_picture(prs.slides[0], IMAGES["slide12_scene"], left=6.35, top=1.28, width=6.25, height=5.45)
    for sh in prs.slides[0].shapes:
        if getattr(sh, "has_text_frame", False) and sh.top > Inches(1.0) and sh.left < Inches(1.0):
            sh.width = Inches(5.25)
            set_font_size(sh, 17, RGBColor(77, 84, 78))
    for page in [5, 7, 14, 15]:
        replace_picture(prs.slides[page - 1], IMAGES["group_graph"], left=6.15, top=1.18, width=6.72, height=5.72)
    replace_picture(prs.slides[15], IMAGES["pair_graph"], left=6.15, top=1.18, width=6.72, height=5.72)
    replace_picture(prs.slides[16], IMAGES["reflection"], left=6.15, top=1.18, width=6.72, height=5.72)
    replace_picture(prs.slides[8], IMAGES["older_teacher_board"], left=6.15, top=1.18, width=6.72, height=5.72)
    replace_picture(prs.slides[9], IMAGES["older_teacher_group"], left=6.15, top=1.18, width=6.72, height=5.72)

    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    main()
