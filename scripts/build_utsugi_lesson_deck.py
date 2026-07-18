from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "output"
PPTX = OUT / "簡報" / "宇津木台語文寫作_20141112_課例分析_可編輯.pptx"
ILLUSTRATIONS = OUT / "圖片" / "宇津木台語文寫作_20141112_GPT插圖"
SCREENSHOTS = OUT / "圖片" / "宇津木台語文寫作_20141112_簡報截圖"

GREEN = RGBColor(32, 70, 48)
CREAM = RGBColor(246, 248, 240)
INK = RGBColor(31, 43, 36)
GOLD = RGBColor(203, 151, 56)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(92, 111, 94)

SLIDES = [
    ("從無字繪本到共同寫作", "材料、同儕與自我，共同構成寫作學習。", ["四段影片合併為 49 分 49 秒完整課堂", "以學生看見、說出、修訂的過程為分析核心", "NotebookLM 筆記提供描述、詮釋、反思架構"], "00:03:00"),
    ("先共同看見，才有共同語言", "共同材料讓學生有可以彼此回應的起點。", ["教師呈現無字繪本", "學生注意顏色與形狀", "圖像成為可共同談論的材料"], "00:04:20"),
    ("無字不是空白，而是敘事的入口", "沒有標準答案，才會出現不同的故事線。", ["學生由觀看展開想像", "每個人可選擇不同角色與事件", "寫作從圖像進入個人經驗"], "00:07:00"),
    ("個人創作保留作者的位置", "先讓每位學生真正開始自己的文本。", ["學生在大紙面書寫或構圖", "安靜書寫是可被尊重的思考時間", "同桌仍保有相互看見的距離"], "00:10:00"),
    ("停筆與觀看也是學習資料", "猶豫不必等於不投入，它也可能是思考正在發生。", ["有人抬頭觀察同伴", "有人停筆後再回到紙面", "觀課先記錄行動，再進行詮釋"], "00:13:00"),
    ("個人工作轉向協同學習", "作品開始流動，寫作成為同儕對話的媒介。", ["學生靠近同伴閱讀作品", "紙面與視線在桌間移動", "成果不只交給教師評定"], "00:18:00"),
    ("傾聽帶來結構再建構", "看見不同表達後，學生回到自己的文本重新安排。", ["同儕作品提供新的敘事可能", "修訂不是尋找教師標準答案", "差異成為再思考的觸媒"], "00:30:10"),
    ("回饋權交給真實讀者", "學生對作者說出感受，讓作品進入公共閱讀。", ["小組內傳閱故事", "以一句話回應作者", "讀者也必須為感受找出語言"], "00:35:45"),
    ("學生相互說話，而非等待判定", "評價權不只在教師，作者能聽見同儕如何閱讀。", ["學生描述同學作品的表現方式", "作者收到真實讀者的回應", "對話從教師中心轉向同儕之間"], "00:38:40"),
    ("困難在共同工作中繼續前進", "學生可以暫時做不出來，仍然留在學習之中。", ["起初不確定能否獨自完成", "同伴與作品提供可依靠的支點", "完成不等於過程中沒有困難"], "00:37:35"),
    ("課堂需要組織條件支持", "一節好課的條件，不只來自一位教師。", ["四人桌讓作品與目光能夠流動", "觀課文化與行政支持同樣重要", "實踐能否持續才是關鍵"], "00:24:00"),
    ("觀課者下一步：追蹤改寫", "從學生的回應學習，而不是急著給教師建議。", ["記錄觀看、對話、回寫的時間點", "比對作品前後稿的實際改變", "以『我從學生的回應學到』進行反思"], "00:45:00"),
]


def text_box(slide, left, top, width, height, text, size, color=INK, bold=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = frame.margin_right = Inches(0)
    frame.margin_top = frame.margin_bottom = Inches(0)
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft JhengHei"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_shape(slide, shape_type, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for index, (title, claim, bullets, timestamp) in enumerate(SLIDES, 1):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = CREAM
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333333, 0.56, GREEN)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 7.32, 13.333333, 0.18, GREEN)
        text_box(slide, 0.46, 0.15, 0.5, 0.22, f"{index:02d}", 18, WHITE)
        text_box(slide, 1.12, 0.17, 5.8, 0.2, "20141112 宇津木台語文寫作｜學習共同體課例分析", 10, WHITE)

        text_box(slide, 0.48, 0.84, 5.95, 0.44, title, 29, GREEN, True)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.48, 1.48, 5.78, 1.08, WHITE, RGBColor(210, 221, 207))
        text_box(slide, 0.72, 1.74, 5.3, 0.54, claim, 23, INK, True)

        y = 2.92
        for bullet in bullets:
            add_shape(slide, MSO_SHAPE.OVAL, 0.7, y + 0.09, 0.12, 0.12, GOLD)
            text_box(slide, 0.98, y, 5.12, 0.48, bullet, 19, INK)
            y += 0.72

        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 6.64, 0.92, 6.22, 4.40, WHITE)
        slide.shapes.add_picture(str(ILLUSTRATIONS / f"slide_{index:02d}.png"), Inches(6.74), Inches(1.02), width=Inches(6.02), height=Inches(3.39))
        text_box(slide, 6.76, 4.57, 2.6, 0.25, "GPT Image 插圖", 12, MUTED)
        text_box(slide, 10.55, 4.57, 2.18, 0.25, f"影片時間：{timestamp}", 12, GREEN)

        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 6.74, 5.08, 2.04, 1.15, WHITE, RGBColor(210, 221, 207))
        slide.shapes.add_picture(str(SCREENSHOTS / f"slide_{index:02d}.png"), Inches(6.80), Inches(5.14), width=Inches(1.92), height=Inches(1.08))
        text_box(slide, 8.98, 5.22, 3.62, 0.3, "原始觀課畫面：以 1080p 合併影片截取", 13, MUTED)
        text_box(slide, 8.98, 5.62, 3.62, 0.38, "文字、眼神、手勢與作品流動，均回看時間碼核對。", 13, INK)
        text_box(slide, 0.48, 7.04, 4.3, 0.14, "NotebookLM《學習共同體哲學文獻重點》分析架構", 8, WHITE)

    PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX)
    print(PPTX)


if __name__ == "__main__":
    build()
