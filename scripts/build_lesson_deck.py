from pathlib import Path
import base64
import io
import textwrap

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PRESENTATION_DIR = ROOT / "output" / "簡報"
IMAGE_DIR = ROOT / "output" / "圖片"
ILLUSTRATION_DIR = IMAGE_DIR / "GPT插圖"
SLIDES_DIR = IMAGE_DIR / "投影片圖片"
PPTX = PRESENTATION_DIR / "hamanosato_fire_lesson_analysis_editable.pptx"
HTML = PRESENTATION_DIR / "hamanosato_fire_lesson_analysis.html"

W, H = 1920, 1080
GREEN = (35, 65, 45)
MINT = (238, 244, 236)
INK = (28, 35, 32)
GOLD = (198, 143, 57)
WHITE = (255, 255, 255)


def font(size, bold=False):
    candidates = [
        "C:/Windows/Fonts/msjhbd.ttc" if bold else "C:/Windows/Fonts/msjh.ttc",
        "C:/Windows/Fonts/NotoSansCJK-Regular.ttc",
        "C:/Windows/Fonts/meiryo.ttc",
    ]
    for path in candidates:
        if Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


F_TITLE = font(58, True)
F_SUB = font(34, False)
F_HEAD = font(37, True)
F_BODY = font(32, False)
F_SMALL = font(24, False)


def wrap(text, width):
    lines = []
    for paragraph in text.split("\n"):
        if " " in paragraph:
            lines.extend(textwrap.wrap(paragraph, width=width, break_long_words=False) or [""])
        else:
            lines.extend([paragraph[i : i + width] for i in range(0, len(paragraph), width)] or [""])
    return lines


slides = [
    {
        "title": "三年級社會：消防與地域公共性",
        "kicker": "2026 日本學習共同體年會・濱之鄉小學實踐",
        "claim": "這堂課把消防從制度知識，推進到「誰來承擔地域安全」的公共性問題。",
        "bullets": [
            "主題不是背誦消防署功能，而是理解消防署、消防團、家庭與地域的互相依存。",
            "核心張力：消防團很重要，但我的父母真的要加入嗎？",
            "學生在情感、危險感、資料證據與公共責任之間來回思考。",
        ],
        "image": "contact_sheet.jpg",
        "time": "00:00-55:46",
    },
    {
        "title": "課前脈絡：教師把單元重新定位",
        "claim": "消防單元被放進公共性與地域參與，而不是只放在機關介紹。",
        "bullets": [
            "教師指出茅崎市具有高火災延燒風險，消防署不是唯一防線。",
            "訪談消防團員，讓教材帶有真實人物、地域歷史與生活脈絡。",
            "課題設計從「地域を守る」推進到「我與家人如何被捲入地域守護」。",
        ],
        "image": "關鍵畫面/00-21-30.jpg",
        "time": "03:20-21:20",
    },
    {
        "title": "本時課題：重要但不想承擔",
        "claim": "孩子遭遇的不是知識不足，而是公共責任與家庭情感的衝突。",
        "bullets": [
            "多數孩子知道消防團重要，卻不希望父母加入。",
            "理由包含危險、忙碌、照顧家庭、父母角色不同。",
            "這份抗拒不是錯誤答案，而是學生真實生活經驗進入社會課。",
        ],
        "image": "關鍵畫面/00-23-30.jpg",
        "time": "22:20-24:30",
    },
    {
        "title": "學習事實 1：回到資料找根據",
        "claim": "學生不是憑印象說，而是被引導回看消防團人數、位置與訪談資料。",
        "bullets": [
            "教師追問哪一個分團、幾個人、離火災地點近不近。",
            "學生重新連結前面學過的消防團資料。",
            "社會課的資料在此成為討論公共責任的依據。",
        ],
        "image": "關鍵畫面/00-27-30.jpg",
        "time": "27:00-30:00",
    },
    {
        "title": "學習事實 2：家長問卷帶來落差",
        "claim": "孩子發現大人也覺得消防團必要，但被邀請時未必會加入。",
        "bullets": [
            "問卷呈現「知道不多」「必要感高」「入團意願低」的分布。",
            "孩子開始讀圖表，不只是聽老師講道理。",
            "期待大人自然承擔的想法被資料動搖。",
        ],
        "image": "關鍵畫面/00-36-30.jpg",
        "time": "36:00-39:00",
    },
    {
        "title": "學習事實 3：讀家長理由",
        "claim": "家長的拒絕不是自私，而是時間、工作、照顧與壓力的現實。",
        "bullets": [
            "學生讀到「一有事就要出動」「有壓力」「工作與家庭忙」等理由。",
            "公共責任變成具體生活負擔，而非抽象口號。",
            "孩子需要同時理解消防團需要與家人限制。",
        ],
        "image": "關鍵畫面/00-39-30.jpg",
        "time": "39:00-42:00",
    },
    {
        "title": "跳躍課題：所有人都加入，如何？",
        "claim": "教師用極端命題讓學生重新檢視自己的立場。",
        "bullets": [
            "如果消防團缺員，是否所有人都應被要求加入？",
            "這不是要學生答應，而是迫使他們處理公共需求與個人限制。",
            "高品質課題讓孩子留在矛盾中，不急著收束成標準答案。",
        ],
        "image": "關鍵畫面/00-43-30.jpg",
        "time": "42:00-46:00",
    },
    {
        "title": "關鍵矛盾：沒人加入會怎樣？",
        "claim": "缺員與再成立的故事，把「不想加入」推回地域安全的現實。",
        "bullets": [
            "教師提出某消防團曾人數歸零、後來由長者重新成立的案例。",
            "學生面對：不想讓家人加入，但無人加入又會使地域失去防線。",
            "這一刻的猶豫，是公共性學習正在發生的證據。",
        ],
        "image": "關鍵畫面/00-46-30.jpg",
        "time": "46:00-49:00",
    },
    {
        "title": "學生的曖昧不是失敗",
        "claim": "學生尚未說清楚答案，反而顯示問題真的進入了他們的生活世界。",
        "bullets": [
            "孩子能感到消防知識與訓練必要，也能感到全員強制不合理。",
            "他們正在從「我家不要」移動到「那地域怎麼辦」。",
            "觀課時要看這種語言的搖擺，而不只看最後結論。",
        ],
        "image": "關鍵畫面/00-51-30.jpg",
        "time": "49:00-52:10",
    },
    {
        "title": "學習共同體詮釋",
        "claim": "本課的價值在於讓學生彼此聽見矛盾，而非被教師快速說服。",
        "bullets": [
            "教材：教科書、消防團員、家長問卷、地域風險共同構成。",
            "關係：學生把父母、消防團員、地域居民放進同一問題。",
            "時間：課堂保留思考的來回，讓公共責任慢慢成形。",
        ],
        "image": "contact_sheet.jpg",
        "time": "全課",
    },
    {
        "title": "觀課者的反思",
        "claim": "公共性教育可以從孩子最真實的家庭感受開始。",
        "bullets": [
            "我學到：孩子的抗拒不是要被修正，而是可被傾聽與深化的材料。",
            "我學到：高品質的社會課會讓學生遭遇「重要但困難」的現實。",
            "我學到：猶豫、停頓、改口，常常比流利正答更接近學習。",
        ],
        "image": "關鍵畫面/00-39-30.jpg",
        "time": "反思",
    },
    {
        "title": "後續可追問的觀課焦點",
        "claim": "下一輪分析可更細看同儕聽取如何改變學生語言。",
        "bullets": [
            "哪些發言讓同伴重新看資料，而不是只表態？",
            "學生如何把「入不入團」改寫成其他地域參與方式？",
            "後續「子ども消防団」活動如何承接本時留下的矛盾？",
        ],
        "image": "關鍵畫面/00-21-30.jpg",
        "time": "延伸",
    },
]

# Each claim is deliberately edited into two complete, presentation-sized lines.
# This avoids orphaned punctuation and accidental phrase splits in both outputs.
CLAIM_LINES = [
    ("消防不只是制度知識，", "更是誰承擔地域安全的公共問題。"),
    ("消防單元放進地域公共性，", "而非只介紹機關。"),
    ("孩子面對的不是知識不足，", "而是責任與家庭情感的衝突。"),
    ("學生不是憑印象說，", "而是回看消防團資料與訪談。"),
    ("孩子知道消防團重要，", "但受邀時未必願意加入。"),
    ("家長拒絕不是自私，", "而是時間、工作與照顧的現實。"),
    ("極端命題讓學生重新", "檢視自己的立場。"),
    ("缺員與再成立的故事，", "讓加入問題回到地域安全。"),
    ("學生尚未說清楚答案，", "正顯示問題進入生活世界。"),
    ("本課讓學生彼此聽見矛盾，", "而非被教師快速說服。"),
    ("公共性教育可從孩子", "真實的家庭感受開始。"),
    ("下一輪可細看同儕聽取", "如何改變學生的語言。"),
]


def claim_lines_for(idx):
    lines = CLAIM_LINES[idx - 1]
    if len(lines) != 2 or any(len(line) > 15 for line in lines):
        raise ValueError(f"Invalid two-line claim for slide {idx}: {lines}")
    return lines


def paste_image(canvas, rel, box):
    path = Path(rel)
    if not path.is_absolute():
        path = IMAGE_DIR / rel
    img = Image.open(path).convert("RGB")
    x, y, w, h = box
    ratio = img.width / img.height
    box_ratio = w / h
    if ratio > box_ratio:
        new_w = w
        new_h = int(w / ratio)
    else:
        new_h = h
        new_w = int(h * ratio)
    img = img.resize((new_w, new_h), Image.LANCZOS)
    left = x + (w - new_w) // 2
    top = y + (h - new_h) // 2
    canvas.paste(img, (left, top))


def draw_slide(data, idx):
    im = Image.new("RGB", (W, H), MINT)
    d = ImageDraw.Draw(im)
    d.rectangle((0, 0, W, 94), fill=GREEN)
    d.rectangle((0, H - 30, W, H), fill=GREEN)
    d.text((72, 24), f"{idx:02d}", fill=(226, 234, 220), font=F_SUB)
    if data.get("kicker"):
        d.text((170, 31), data["kicker"], fill=(226, 234, 220), font=F_SMALL)
    d.text((72, 122), data["title"], fill=GREEN, font=F_TITLE)
    d.rounded_rectangle((74, 210, 940, 398), radius=18, fill=WHITE, outline=(211, 222, 207), width=2)
    claim_lines = claim_lines_for(idx)
    d.text((108, 248), "\n".join(claim_lines), fill=INK, font=F_HEAD, spacing=10)
    y = 440
    for bullet in data["bullets"]:
        d.ellipse((102, y + 15, 120, y + 33), fill=GOLD)
        for line in wrap(bullet, 20):
            d.text((142, y), line, fill=INK, font=F_BODY)
            y += 44
        y += 18
    # Keep a visible gutter between the enlarged text column and illustration.
    img_box = (1010, 185, 840, 560)
    d.rounded_rectangle((img_box[0] - 10, img_box[1] - 10, img_box[0] + img_box[2] + 10, img_box[1] + img_box[3] + 10), radius=18, fill=WHITE)
    paste_image(im, ILLUSTRATION_DIR / f"slide_{idx:02d}.png", img_box)
    d.text((1010, 778), f"影片時間：{data['time']}", fill=GREEN, font=F_SUB)
    d.text((72, 1015), "和北極星境遇｜學習共同體公開課影片分析", fill=(226, 234, 220), font=F_SMALL)
    return im


def add_ppt_textbox(slide, left, top, width, height, text, size, bold=False, color=INK):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = "Microsoft JhengHei"
    p.font.color.rgb = RGBColor(*color)
    return box


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for idx, data in enumerate(slides, 1):
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(*MINT)

        header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.58))
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(*GREEN)
        header.line.fill.background()
        footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.32), prs.slide_width, Inches(0.18))
        footer.fill.solid()
        footer.fill.fore_color.rgb = RGBColor(*GREEN)
        footer.line.fill.background()

        add_ppt_textbox(slide, 0.48, 0.14, 0.5, 0.28, f"{idx:02d}", 20, False, (226, 234, 220))
        if data.get("kicker"):
            add_ppt_textbox(slide, 1.15, 0.17, 5.7, 0.22, data["kicker"], 10, False, (226, 234, 220))
        # Long Chinese titles are deliberately wrapped before the claim card so
        # they never intrude into its bordered area.
        title_lines = wrap(data["title"], 17)
        title_size = 31 if len(title_lines) == 1 and len(data["title"]) <= 14 else 27
        if len(title_lines) > 1:
            title_size = 25
        title_height = len(title_lines) * (0.42 if len(title_lines) == 1 else 0.37)
        add_ppt_textbox(slide, 0.5, 0.86, 6.1, title_height, "\n".join(title_lines), title_size, True, GREEN)

        claim_lines = claim_lines_for(idx)
        # PowerPoint Chinese glyphs need more leading than the nominal point size suggests.
        claim_height = 0.42 + len(claim_lines) * 0.54
        claim_top = max(1.5, 0.86 + title_height + 0.22)
        claim_card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.48),
            Inches(claim_top),
            Inches(6.1),
            Inches(claim_height),
        )
        claim_card.fill.solid()
        claim_card.fill.fore_color.rgb = RGBColor(*WHITE)
        claim_card.line.color.rgb = RGBColor(211, 222, 207)
        add_ppt_textbox(slide, 0.72, claim_top + 0.19, 5.62, claim_height - 0.3, "\n".join(claim_lines), 22, True, INK)

        bullet_y = claim_top + claim_height + 0.35
        for bullet in data["bullets"]:
            bullet_lines = wrap(bullet, 19)
            bullet_height = len(bullet_lines) * 0.38
            dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.7), Inches(bullet_y + 0.08), Inches(0.13), Inches(0.13))
            dot.fill.solid()
            dot.fill.fore_color.rgb = RGBColor(*GOLD)
            dot.line.fill.background()
            add_ppt_textbox(slide, 0.98, bullet_y, 5.35, bullet_height + 0.04, "\n".join(bullet_lines), 20, False, INK)
            bullet_y += bullet_height + 0.18

        image = ILLUSTRATION_DIR / f"slide_{idx:02d}.png"
        frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(1.34), Inches(5.92), Inches(3.52))
        frame.fill.solid()
        frame.fill.fore_color.rgb = RGBColor(*WHITE)
        frame.line.fill.background()
        slide.shapes.add_picture(str(image), Inches(6.95), Inches(1.44), width=Inches(5.72), height=Inches(3.22))
        add_ppt_textbox(slide, 6.95, 5.02, 4.6, 0.3, f"影片時間：{data['time']}", 17, False, GREEN)
        add_ppt_textbox(slide, 0.5, 7.05, 4.0, 0.16, "和北極星境遇｜學習共同體公開課影片分析", 8, False, (226, 234, 220))
    PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX)


def build_html():
    sections = []
    for idx, path in enumerate(sorted(SLIDES_DIR.glob("slide_*.png")), 1):
        with Image.open(path) as im:
            buf = io.BytesIO()
            im.save(buf, format="WEBP", quality=88, method=6)
        data = base64.b64encode(buf.getvalue()).decode("ascii")
        active = " active" if idx == 1 else ""
        sections.append(f'<section class="slide{active}" aria-hidden="{str(idx != 1).lower()}"><img src="data:image/webp;base64,{data}" alt="Slide {idx}"></section>')
    html = f"""<!doctype html>
<html lang="zh-TW">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>濱之鄉消防課分析</title>
<style>
*{{box-sizing:border-box}}
html,body{{margin:0;width:100%;height:100%;background:#eef4ec;font-family:"Microsoft JhengHei",sans-serif;overflow:hidden}}
body{{display:grid;place-items:center;padding:8px}}
.shell{{display:block;width:fit-content;max-width:98vw}}
.deck{{height:min(98vh,calc(98vw * 0.5625),888px);aspect-ratio:16/9;background:#fff;box-shadow:0 12px 32px #1935242e}}
.stage{{width:100%;height:100%;display:grid;place-items:center}}
.slide{{display:none;width:100%;height:100%;align-items:center;justify-content:center}}
.slide.active{{display:flex}}
.slide img{{display:block;width:100%;height:100%;object-fit:contain}}
@media(max-width:720px){{body{{padding:0}}.shell{{max-width:100vw;width:100vw}}.deck{{height:min(100vh,calc(100vw * 0.5625));width:auto}}}}
@media print{{html,body{{background:#fff;overflow:visible}}body{{display:block;padding:0}}.shell,.deck{{width:100vw;height:100vh;max-width:none;display:block;box-shadow:none;background:#fff}}.stage,.slide,.slide img{{width:100vw;height:100vh;max-width:none;max-height:none}}}}
</style>
</head>
<body>
<div class="shell"><main class="deck"><div class="stage">{''.join(sections)}</div></main></div>
<script>
const slides=[...document.querySelectorAll('.slide')];let current=0;function render(){{slides.forEach((s,i)=>{{const a=i===current;s.classList.toggle('active',a);s.setAttribute('aria-hidden',String(!a));}});}}function move(d){{current=Math.max(0,Math.min(slides.length-1,current+d));render();}}document.addEventListener('keydown',e=>{{if(e.key==='ArrowLeft')move(-1);if(e.key==='ArrowRight'||e.key===' '){{e.preventDefault();move(1);}}}});render();
</script>
</body>
</html>"""
    HTML.write_text(html, encoding="utf-8")


def main():
    PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)
    IMAGE_DIR.mkdir(parents=True, exist_ok=True)
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)
    for idx, data in enumerate(slides, 1):
        draw_slide(data, idx).save(SLIDES_DIR / f"slide_{idx:02d}.png")
    build_pptx()
    build_html()
    print(PPTX)
    print(HTML)
    print(SLIDES_DIR)


if __name__ == "__main__":
    main()
